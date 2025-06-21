#!/usr/bin/env python3
"""
Bamini to Unicode Tamil PowerPoint Converter
Converts Bamini encoded Tamil text to Unicode Tamil in PowerPoint files
while preserving all formatting and background images.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
import argparse
import collections

# Bamini25 to Unicode Tamil character mapping
# Based on actual Bamini25 font character codes
BAMINI_TO_UNICODE = collections.OrderedDict()

BAMINI_TO_UNICODE["sp"] = "ளி"
BAMINI_TO_UNICODE["hp"] = "ரி"
BAMINI_TO_UNICODE["hP"] = "ரீ"
BAMINI_TO_UNICODE["uP"] = "ரீ"
BAMINI_TO_UNICODE["u;"] = "ர்"
BAMINI_TO_UNICODE["h;"] = "ர்"
BAMINI_TO_UNICODE["H"] = "ர்"
BAMINI_TO_UNICODE["nfs"] = "கௌ"
BAMINI_TO_UNICODE["Nfh"] = "கோ"
BAMINI_TO_UNICODE["nfh"] = "கொ"
BAMINI_TO_UNICODE["fh"] = "கா"
BAMINI_TO_UNICODE["fp"] = "கி"
BAMINI_TO_UNICODE["fP"] = "கீ"
BAMINI_TO_UNICODE["F"] = "கு"
BAMINI_TO_UNICODE["$"] = "கூ"
BAMINI_TO_UNICODE["nf"] = "கெ"
BAMINI_TO_UNICODE["Nf"] = "கே"
BAMINI_TO_UNICODE["if"] = "கை"
BAMINI_TO_UNICODE["f;"] = "க்"
BAMINI_TO_UNICODE["f"] = "க"
BAMINI_TO_UNICODE["nqs"] = "ஙௌ"
BAMINI_TO_UNICODE["Nqh"] = "ஙோ"
BAMINI_TO_UNICODE["nqh"] = "ஙொ"
BAMINI_TO_UNICODE["qh"] = "ஙா"
BAMINI_TO_UNICODE["qp"] = "ஙி"
BAMINI_TO_UNICODE["qP"] = "ஙீ"
BAMINI_TO_UNICODE["nq"] = "ஙெ"
BAMINI_TO_UNICODE["Nq"] = "ஙே"
BAMINI_TO_UNICODE["iq"] = "ஙை"
BAMINI_TO_UNICODE["q;"] = "ங்"
BAMINI_TO_UNICODE["q"] = "ங"
BAMINI_TO_UNICODE["nrs"] = "சௌ"
BAMINI_TO_UNICODE["Nrh"] = "சோ"
BAMINI_TO_UNICODE["nrh"] = "சொ"
BAMINI_TO_UNICODE["rh"] = "சா"
BAMINI_TO_UNICODE["rp"] = "சி"
BAMINI_TO_UNICODE["rP"] = "சீ"
BAMINI_TO_UNICODE["R"] = "சு"
BAMINI_TO_UNICODE["#"] = "சூ"
BAMINI_TO_UNICODE["nr"] = "செ"
BAMINI_TO_UNICODE["Nr"] = "சே"
BAMINI_TO_UNICODE["ir"] = "சை"
BAMINI_TO_UNICODE["r;"] = "ச்"
BAMINI_TO_UNICODE["r"] = "ச"
BAMINI_TO_UNICODE["n[s"] = "ஜௌ"
BAMINI_TO_UNICODE["N[h"] = "ஜோ"
BAMINI_TO_UNICODE["n[h"] = "ஜொ"
BAMINI_TO_UNICODE["[h"] = "ஜா"
BAMINI_TO_UNICODE["[p"] = "ஜி"
BAMINI_TO_UNICODE["[P"] = "ஜீ"
BAMINI_TO_UNICODE["[{"] = "ஜு"
BAMINI_TO_UNICODE["[_"] = "ஜூ"
BAMINI_TO_UNICODE["n["] = "ஜெ"
BAMINI_TO_UNICODE["N["] = "ஜே"
BAMINI_TO_UNICODE["i["] = "ஜை"
BAMINI_TO_UNICODE["[;"] = "ஜ்"
BAMINI_TO_UNICODE["["] = "ஜ"
BAMINI_TO_UNICODE["nQs"] = "ஞௌ"
BAMINI_TO_UNICODE["NQh"] = "ஞோ"
BAMINI_TO_UNICODE["nQh"] = "ஞொ"
BAMINI_TO_UNICODE["Qh"] = "ஞா"
BAMINI_TO_UNICODE["Qp"] = "ஞி"
BAMINI_TO_UNICODE["QP"] = "ஞீ"
BAMINI_TO_UNICODE["nQ"] = "ஞெ"
BAMINI_TO_UNICODE["NQ"] = "ஞே"
BAMINI_TO_UNICODE["iQ"] = "ஞை"
BAMINI_TO_UNICODE["Q;"] = "ஞ்"
BAMINI_TO_UNICODE["Q"] = "ஞ"
BAMINI_TO_UNICODE["nls"] = "டௌ"
BAMINI_TO_UNICODE["Nlh"] = "டோ"
BAMINI_TO_UNICODE["nlh"] = "டொ"
BAMINI_TO_UNICODE["lp"] = "டி"
BAMINI_TO_UNICODE["lP"] = "டீ"
BAMINI_TO_UNICODE["lh"] = "டா"
BAMINI_TO_UNICODE["b"] = "டி"
BAMINI_TO_UNICODE["B"] = "டீ"
BAMINI_TO_UNICODE["L"] = "டு"
BAMINI_TO_UNICODE["^"] = "டூ"
BAMINI_TO_UNICODE["nl"] = "டெ"
BAMINI_TO_UNICODE["Nl"] = "டே"
BAMINI_TO_UNICODE["il"] = "டை"
BAMINI_TO_UNICODE["l;"] = "ட்"
BAMINI_TO_UNICODE["l"] = "ட"
BAMINI_TO_UNICODE["nzs"] = "ணௌ"
BAMINI_TO_UNICODE["Nzh"] = "ணோ"
BAMINI_TO_UNICODE["nzh"] = "ணொ"
BAMINI_TO_UNICODE["zh"] = "ணா"
BAMINI_TO_UNICODE["zp"] = "ணி"
BAMINI_TO_UNICODE["zP"] = "ணீ"
BAMINI_TO_UNICODE["Zh"] = "ணூ"
BAMINI_TO_UNICODE["Z}"] = "ணூ"
BAMINI_TO_UNICODE["nz"] = "ணெ"
BAMINI_TO_UNICODE["Nz"] = "ணே"
BAMINI_TO_UNICODE["iz"] = "ணை"
BAMINI_TO_UNICODE["z;"] = "ண்"
BAMINI_TO_UNICODE["Z"] = "ணு"
BAMINI_TO_UNICODE["z"] = "ண"
BAMINI_TO_UNICODE["njs"] = "தௌ"
BAMINI_TO_UNICODE["Njh"] = "தோ"
BAMINI_TO_UNICODE["njh"] = "தொ"
BAMINI_TO_UNICODE["jh"] = "தா"
BAMINI_TO_UNICODE["jp"] = "தி"
BAMINI_TO_UNICODE["jP"] = "தீ"
BAMINI_TO_UNICODE["Jh"] = "தூ"
BAMINI_TO_UNICODE["Jh"] = "தூ"
BAMINI_TO_UNICODE["J}"] = "தூ"
BAMINI_TO_UNICODE["J"] = "து"
BAMINI_TO_UNICODE["nj"] = "தெ"
BAMINI_TO_UNICODE["Nj"] = "தே"
BAMINI_TO_UNICODE["ij"] = "தை"
BAMINI_TO_UNICODE["j;"] = "த்"
BAMINI_TO_UNICODE["j"] = "த"
BAMINI_TO_UNICODE["nes"] = "நௌ"
BAMINI_TO_UNICODE["Neh"] = "நோ"
BAMINI_TO_UNICODE["neh"] = "நொ"
BAMINI_TO_UNICODE["eh"] = "நா"
BAMINI_TO_UNICODE["ep"] = "நி"
BAMINI_TO_UNICODE["eP"] = "நீ"
BAMINI_TO_UNICODE["E}"] = "நூ"
BAMINI_TO_UNICODE["Eh"] = "நூ"
BAMINI_TO_UNICODE["E"] = "நு"
BAMINI_TO_UNICODE["ne"] = "நெ"
BAMINI_TO_UNICODE["Ne"] = "நே"
BAMINI_TO_UNICODE["ie"] = "நை"
BAMINI_TO_UNICODE["e;"] = "ந்"
BAMINI_TO_UNICODE["e"] = "ந"
BAMINI_TO_UNICODE["nds"] = "னௌ"
BAMINI_TO_UNICODE["Ndh"] = "னோ"
BAMINI_TO_UNICODE["ndh"] = "னொ"
BAMINI_TO_UNICODE["dh"] = "னா"
BAMINI_TO_UNICODE["dp"] = "னி"
BAMINI_TO_UNICODE["dP"] = "னீ"
BAMINI_TO_UNICODE["D}"] = "னூ"
BAMINI_TO_UNICODE["Dh"] = "னூ"
BAMINI_TO_UNICODE["D"] = "னு"
BAMINI_TO_UNICODE["nd"] = "னெ"
BAMINI_TO_UNICODE["Nd"] = "னே"
BAMINI_TO_UNICODE["id"] = "னை"
BAMINI_TO_UNICODE["d;"] = "ன்"
BAMINI_TO_UNICODE["d"] = "ன"
BAMINI_TO_UNICODE["ngs"] = "பௌ"
BAMINI_TO_UNICODE["Ngh"] = "போ"
BAMINI_TO_UNICODE["ngh"] = "பொ"
BAMINI_TO_UNICODE["gh"] = "பா"
BAMINI_TO_UNICODE["gp"] = "பி"
BAMINI_TO_UNICODE["gP"] = "பீ"
BAMINI_TO_UNICODE["G"] = "பு"
BAMINI_TO_UNICODE["ng"] = "பெ"
BAMINI_TO_UNICODE["Ng"] = "பே"
BAMINI_TO_UNICODE["ig"] = "பை"
BAMINI_TO_UNICODE["g;"] = "ப்"
BAMINI_TO_UNICODE["g"] = "ப"
BAMINI_TO_UNICODE["nks"] = "மௌ"
BAMINI_TO_UNICODE["Nkh"] = "மோ"
BAMINI_TO_UNICODE["nkh"] = "மொ"
BAMINI_TO_UNICODE["kh"] = "மா"
BAMINI_TO_UNICODE["kp"] = "மி"
BAMINI_TO_UNICODE["kP"] = "மீ"
BAMINI_TO_UNICODE["K"] = "மு"
BAMINI_TO_UNICODE["%"] = "மூ"
BAMINI_TO_UNICODE["nk"] = "மெ"
BAMINI_TO_UNICODE["Nk"] = "மே"
BAMINI_TO_UNICODE["ik"] = "மை"
BAMINI_TO_UNICODE["k;"] = "ம்"
BAMINI_TO_UNICODE["k"] = "ம"
BAMINI_TO_UNICODE["nas"] = "யௌ"
BAMINI_TO_UNICODE["Nah"] = "யோ"
BAMINI_TO_UNICODE["nah"] = "யொ"
BAMINI_TO_UNICODE["ah"] = "யா"
BAMINI_TO_UNICODE["ap"] = "யி"
BAMINI_TO_UNICODE["aP"] = "யீ"
BAMINI_TO_UNICODE["A"] = "யு"
BAMINI_TO_UNICODE["A+"] = "யூ"
BAMINI_TO_UNICODE["na"] = "யெ"
BAMINI_TO_UNICODE["Na"] = "யே"
BAMINI_TO_UNICODE["ia"] = "யை"
BAMINI_TO_UNICODE["a;"] = "ய்"
BAMINI_TO_UNICODE["a"] = "ய"
BAMINI_TO_UNICODE["nus"] = "ரௌ"
BAMINI_TO_UNICODE["Nuh"] = "ரோ"
BAMINI_TO_UNICODE["nuh"] = "ரொ"
BAMINI_TO_UNICODE["uh"] = "ரா"
BAMINI_TO_UNICODE["up"] = "ரி"
BAMINI_TO_UNICODE["U"] = "ரு"
BAMINI_TO_UNICODE["&"] = "ரூ"
BAMINI_TO_UNICODE["nu"] = "ரெ"
BAMINI_TO_UNICODE["Nu"] = "ரே"
BAMINI_TO_UNICODE["iu"] = "ரை"
BAMINI_TO_UNICODE["u"] = "ர"
BAMINI_TO_UNICODE["nys"] = "லௌ"
BAMINI_TO_UNICODE["Nyh"] = "லோ"
BAMINI_TO_UNICODE["nyh"] = "லொ"
BAMINI_TO_UNICODE["yh"] = "லா"
BAMINI_TO_UNICODE["yp"] = "லி"
BAMINI_TO_UNICODE["yP"] = "லீ"
BAMINI_TO_UNICODE["Yh"] = "லூ"
BAMINI_TO_UNICODE["Y}"] = "லூ"
BAMINI_TO_UNICODE["Y"] = "லு"
BAMINI_TO_UNICODE["ny"] = "லெ"
BAMINI_TO_UNICODE["Ny"] = "லே"
BAMINI_TO_UNICODE["iy"] = "லை"
BAMINI_TO_UNICODE["y;"] = "ல்"
BAMINI_TO_UNICODE["y"] = "ல"
BAMINI_TO_UNICODE["nss"] = "ளௌ"
BAMINI_TO_UNICODE["Nsh"] = "ளோ"
BAMINI_TO_UNICODE["nsh"] = "ளொ"
BAMINI_TO_UNICODE["sh"] = "ளா"
BAMINI_TO_UNICODE["sP"] = "ளீ"
BAMINI_TO_UNICODE["Sh"] = "ளூ"
BAMINI_TO_UNICODE["S"] = "ளு"
BAMINI_TO_UNICODE["ns"] = "ளெ"
BAMINI_TO_UNICODE["Ns"] = "ளே"
BAMINI_TO_UNICODE["is"] = "ளை"
BAMINI_TO_UNICODE["s;"] = "ள்"
BAMINI_TO_UNICODE["s"] = "ள"
BAMINI_TO_UNICODE["ntt"] = "வௌ"
BAMINI_TO_UNICODE["Nth"] = "வோ"
BAMINI_TO_UNICODE["nth"] = "வொ"
BAMINI_TO_UNICODE["th"] = "வா"
BAMINI_TO_UNICODE["tp"] = "வி"
BAMINI_TO_UNICODE["tP"] = "வீ"
BAMINI_TO_UNICODE["nt"] = "வெ"
BAMINI_TO_UNICODE["Nt"] = "வே"
BAMINI_TO_UNICODE["it"] = "வை"
BAMINI_TO_UNICODE["t;"] = "வ்"
BAMINI_TO_UNICODE["t"] = "வ"
BAMINI_TO_UNICODE["noo"] = "ழௌ"
BAMINI_TO_UNICODE["Noh"] = "ழோ"
BAMINI_TO_UNICODE["noh"] = "ழொ"
BAMINI_TO_UNICODE["oh"] = "ழா"
BAMINI_TO_UNICODE["op"] = "ழி"
BAMINI_TO_UNICODE["oP"] = "ழீ"
BAMINI_TO_UNICODE["*"] = "ழூ"
BAMINI_TO_UNICODE["O"] = "ழு"
BAMINI_TO_UNICODE["no"] = "ழெ"
BAMINI_TO_UNICODE["No"] = "ழே"
BAMINI_TO_UNICODE["io"] = "ழை"
BAMINI_TO_UNICODE["o;"] = "ழ்"
BAMINI_TO_UNICODE["o"] = "ழ"
BAMINI_TO_UNICODE["nws"] = "றௌ"
BAMINI_TO_UNICODE["Nwh"] = "றோ"
BAMINI_TO_UNICODE["nwh"] = "றொ"
BAMINI_TO_UNICODE["wh"] = "றா"
BAMINI_TO_UNICODE["wp"] = "றி"
BAMINI_TO_UNICODE["wP"] = "றீ"
BAMINI_TO_UNICODE["Wh"] = "றூ"
BAMINI_TO_UNICODE["W}"] = "றூ"
BAMINI_TO_UNICODE["W"] = "று"
BAMINI_TO_UNICODE["nw"] = "றெ"
BAMINI_TO_UNICODE["Nw"] = "றே"
BAMINI_TO_UNICODE["iw"] = "றை"
BAMINI_TO_UNICODE["w;"] = "ற்"
BAMINI_TO_UNICODE["w"] = "ற"
BAMINI_TO_UNICODE["n``"] = "ஹௌ"
BAMINI_TO_UNICODE["N`h"] = "ஹோ"
BAMINI_TO_UNICODE["n`h"] = "ஹொ"
BAMINI_TO_UNICODE["`h"] = "ஹா"
BAMINI_TO_UNICODE["`p"] = "ஹி"
BAMINI_TO_UNICODE["`P"] = "ஹீ"
BAMINI_TO_UNICODE["n`"] = "ஹெ"
BAMINI_TO_UNICODE["N`"] = "ஹே"
BAMINI_TO_UNICODE["i`"] = "ஹை"
BAMINI_TO_UNICODE["`;"] = "ஹ்"
BAMINI_TO_UNICODE["`"] = "ஹ"
BAMINI_TO_UNICODE["n\s"] = "ஷௌ"
BAMINI_TO_UNICODE["N\h"] = "ஷோ"
BAMINI_TO_UNICODE["n\h"] = "ஷொ"
BAMINI_TO_UNICODE["\h"] = "ஷா"
BAMINI_TO_UNICODE["\p"] = "ஷி"
BAMINI_TO_UNICODE["\P"] = "ஷீ"
BAMINI_TO_UNICODE["n\\"] = "ஷெ"
BAMINI_TO_UNICODE["N\\"] = "ஷே"
BAMINI_TO_UNICODE["i\\"] = "ஷை"
BAMINI_TO_UNICODE["\\;"] = "ஷ்"
BAMINI_TO_UNICODE["\\"] = "ஷ"
BAMINI_TO_UNICODE["n]s"] = "ஸௌ"
BAMINI_TO_UNICODE["N]h"] = "ஸோ"
BAMINI_TO_UNICODE["n]h"] = "ஸொ"
BAMINI_TO_UNICODE["]h"] = "ஸா"
BAMINI_TO_UNICODE["]p"] = "ஸி"
BAMINI_TO_UNICODE["]P"] = "ஸீ"
BAMINI_TO_UNICODE["n]"] = "ஸெ"
BAMINI_TO_UNICODE["N]"] = "ஸே"
BAMINI_TO_UNICODE["i]"] = "ஸை"
BAMINI_TO_UNICODE["];"] = "ஸ்"
BAMINI_TO_UNICODE["]"] = "ஸ"
BAMINI_TO_UNICODE["m"] = "அ"
BAMINI_TO_UNICODE["M"] = "ஆ"
BAMINI_TO_UNICODE["<"] = "ஈ"
BAMINI_TO_UNICODE["c"] = "உ"
BAMINI_TO_UNICODE["C"] = "ஊ"
BAMINI_TO_UNICODE["v"] = "எ"
BAMINI_TO_UNICODE["V"] = "ஏ"
BAMINI_TO_UNICODE["I"] = "ஐ"
BAMINI_TO_UNICODE["x"] = "ஒ"
BAMINI_TO_UNICODE["X"] = "ஓ"
BAMINI_TO_UNICODE["xs"] = "ஔ"
BAMINI_TO_UNICODE["/"] = "ஃ"
BAMINI_TO_UNICODE[","] = "இ"
BAMINI_TO_UNICODE["="] = "ஸ்ரீ"
BAMINI_TO_UNICODE[">"] = ","
BAMINI_TO_UNICODE["T"] = "வு"
BAMINI_TO_UNICODE["வு+"] = "வூ"
BAMINI_TO_UNICODE["பு+"] = "பூ"
BAMINI_TO_UNICODE["யு+"] = "யூ"
BAMINI_TO_UNICODE["சு+"] = "சூ"
BAMINI_TO_UNICODE["+"] = "ooh"
BAMINI_TO_UNICODE[";"] = "்"
BAMINI_TO_UNICODE["@"] = ";"



BAMINI_TO_UNICODE["¿f"] = "கை"
BAMINI_TO_UNICODE["¿q"] = "ஙை"
BAMINI_TO_UNICODE["¿r"] = "சை"
BAMINI_TO_UNICODE["¿["] = "ஜை"
BAMINI_TO_UNICODE["¿Q"] = "ஞை"
BAMINI_TO_UNICODE["¿l"] = "டை"
BAMINI_TO_UNICODE["¿z"] = "ணை"
BAMINI_TO_UNICODE["¿j"] = "தை"
BAMINI_TO_UNICODE["¿e"] = "நை"
BAMINI_TO_UNICODE["¿d"] = "னை"
BAMINI_TO_UNICODE["¿g"] = "பை"
BAMINI_TO_UNICODE["¿k"] = "மை"
BAMINI_TO_UNICODE["¿a"] = "யை"
BAMINI_TO_UNICODE["¿u"] = "ரை"
BAMINI_TO_UNICODE["¿y"] = "லை"
BAMINI_TO_UNICODE["¿s"] = "ளை"
BAMINI_TO_UNICODE["¿t"] = "வை"
BAMINI_TO_UNICODE["¿o"] = "ழை"
BAMINI_TO_UNICODE["¿w"] = "றை"
BAMINI_TO_UNICODE["¿`"] = "ஹை"
BAMINI_TO_UNICODE["¿\\"] = "ஷை"
BAMINI_TO_UNICODE["¿]"] = "ஸை"
BAMINI_TO_UNICODE["¿"] = "ை"
BAMINI_TO_UNICODE["≈"] = "ௐ"

def convert_bamini_to_unicode(text):
    """Convert Bamini encoded text to Unicode Tamil."""
    if not text:
        return text
    
    # Sort mappings by length (longest first) to avoid partial replacements
    sorted_mappings = sorted(BAMINI_TO_UNICODE.items(), key=lambda x: len(x[0]), reverse=True)
    
    converted_text = text
    for bamini_char, unicode_char in sorted_mappings:
        converted_text = converted_text.replace(bamini_char, unicode_char)
    
    return converted_text

def convert_text_in_shape(shape):
    """Convert text in a shape while preserving formatting."""
    if not shape.has_text_frame:
        return
    
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                # Store original formatting
                original_font = run.font
                original_text = run.text
                
                # Convert text
                converted_text = convert_bamini_to_unicode(original_text)
                
                # Update text while preserving formatting
                run.text = converted_text
                
                # Ensure font properties are maintained
                if original_font.name:
                    run.font.name = original_font.name
                if original_font.size:
                    run.font.size = original_font.size
                if original_font.bold is not None:
                    run.font.bold = original_font.bold
                if original_font.italic is not None:
                    run.font.italic = original_font.italic
                if original_font.underline is not None:
                    run.font.underline = original_font.underline
                
                # Handle color carefully to avoid RGB errors
                try:
                    if hasattr(original_font, 'color') and original_font.color:
                        if hasattr(original_font.color, 'rgb') and original_font.color.rgb is not None:
                            run.font.color.rgb = original_font.color.rgb
                        elif hasattr(original_font.color, 'theme_color') and original_font.color.theme_color is not None:
                            run.font.color.theme_color = original_font.color.theme_color
                except Exception:
                    # Skip color if there are issues
                    pass

def convert_text_in_table(table):
    """Convert text in table cells while preserving formatting."""
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text:
                        # Store original formatting
                        original_font = run.font
                        original_text = run.text
                        
                        # Convert text
                        converted_text = convert_bamini_to_unicode(original_text)
                        
                        # Update text while preserving formatting
                        run.text = converted_text
                        
                        # Maintain formatting
                        if original_font.name:
                            run.font.name = original_font.name
                        if original_font.size:
                            run.font.size = original_font.size
                        if original_font.bold is not None:
                            run.font.bold = original_font.bold
                        if original_font.italic is not None:
                            run.font.italic = original_font.italic
                        
                        # Handle color carefully to avoid RGB errors
                        try:
                            if hasattr(original_font, 'color') and original_font.color:
                                if hasattr(original_font.color, 'rgb') and original_font.color.rgb is not None:
                                    run.font.color.rgb = original_font.color.rgb
                                elif hasattr(original_font.color, 'theme_color') and original_font.color.theme_color is not None:
                                    run.font.color.theme_color = original_font.color.theme_color
                        except Exception:
                            # Skip color if there are issues
                            pass

def convert_presentation(input_file, output_file=None):
    """Convert Bamini text to Unicode Tamil in PowerPoint presentation."""
    try:
        # Load presentation
        print(f"Loading presentation: {input_file}")
        presentation = Presentation(input_file)
        
        slide_count = 0
        text_conversions = 0
        
        # Process each slide
        for slide_num, slide in enumerate(presentation.slides, 1):
            print(f"Processing slide {slide_num}...")
            slide_count += 1
            
            # Process all shapes in the slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Count text runs before conversion
                    original_runs = 0
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                original_runs += 1
                    
                    # Convert text in shape
                    convert_text_in_shape(shape)
                    text_conversions += original_runs
                
                # Handle tables separately
                elif hasattr(shape, 'table'):
                    convert_text_in_table(shape.table)
                    # Count table cells
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip():
                                        text_conversions += 1
        
        # Save converted presentation
        if output_file is None:
            # Create output filename
            base_name = os.path.splitext(input_file)[0]
            output_file = f"{base_name}_unicode_tamil.pptx"
        
        print(f"Saving converted presentation: {output_file}")
        presentation.save(output_file)
        
        print(f"\nConversion completed successfully!")
        print(f"Slides processed: {slide_count}")
        print(f"Text elements converted: {text_conversions}")
        print(f"Output saved as: {output_file}")
        
        return True
        
    except FileNotFoundError:
        print(f"Error: Input file '{input_file}' not found.")
        return False
    except Exception as e:
        print(f"Error during conversion: {str(e)}")
        return False

def main():
    parser = argparse.ArgumentParser(
        description="Convert Bamini encoded Tamil text to Unicode Tamil in PowerPoint files",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python bamini_converter.py input.pptx
  python bamini_converter.py input.pptx -o output.pptx
  python bamini_converter.py presentation.pptx --output converted_presentation.pptx
        """
    )
    
    parser.add_argument('input_file', 
                       help='Input PowerPoint file (.pptx) with Bamini text')
    parser.add_argument('-o', '--output', 
                       help='Output file name (optional, defaults to input_unicode_tamil.pptx)')
    
    args = parser.parse_args()
    
    # Validate input file
    if not os.path.exists(args.input_file):
        print(f"Error: Input file '{args.input_file}' does not exist.")
        sys.exit(1)
    
    if not args.input_file.lower().endswith('.pptx'):
        print("Error: Input file must be a PowerPoint file (.pptx)")
        sys.exit(1)
    
    # Convert presentation
    success = convert_presentation(args.input_file, args.output)
    
    if success:
        print("\n✓ Conversion completed successfully!")
        print("All formatting and background images have been preserved.")
    else:
        print("\n✗ Conversion failed!")
        sys.exit(1)

if __name__ == "__main__":
    main()
